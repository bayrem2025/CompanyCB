import os
import pandas as pd
from pathlib import Path
from typing import Optional, Tuple, List
import logging
from langchain.text_splitter import RecursiveCharacterTextSplitter
from config import CHUNK_SIZE, CHUNK_OVERLAP
from langchain.schema import Document

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class UniversalExtractor:
    """Extracteur de texte universel"""
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[Optional[str], dict]:
        """Extrait le texte avec des bibliothèques spécifiques par format"""
        try:
            ext = Path(file_path).suffix.lower()
            metadata = {
                "source": file_path,
                "file_type": ext[1:],
                "file_name": os.path.basename(file_path)
            }

            if ext == '.pdf':
                return UniversalExtractor._extract_pdf(file_path, metadata)
            elif ext in ['.xlsx', '.xls']:
                return UniversalExtractor._extract_excel(file_path, metadata)
            elif ext in ['.docx']:
                return UniversalExtractor._extract_docx(file_path, metadata)
            elif ext in ['.pptx']:
                return UniversalExtractor._extract_pptx(file_path, metadata)
            elif ext == '.csv':
                return UniversalExtractor._extract_csv(file_path, metadata)
            elif ext == '.txt':
                return UniversalExtractor._extract_txt(file_path, metadata)
            else:
                logger.warning(f"Format non supporté : {file_path}")
                return None, {}

        except Exception as e:
            logger.error(f"Erreur avec {file_path}: {str(e)}")
            return None, {}

    @staticmethod
    def _extract_pdf(file_path: str, metadata: dict) -> Tuple[str, dict]:
        """Extrait le texte d'un PDF avec pypdf"""
        from pypdf import PdfReader
        text = []
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            metadata["pages"] = len(reader.pages)
            for i, page in enumerate(reader.pages):
                text.append(f"\n[PAGE {i+1}]\n{page.extract_text()}")
        return "\n".join(text), metadata

    @staticmethod
    def _extract_excel(file_path: str, metadata: dict) -> Tuple[str, dict]:
        """Extrait le texte d'Excel avec pandas"""
        text = []
        with pd.ExcelFile(file_path) as xls:
            metadata["sheets"] = xls.sheet_names
            for sheet_name in xls.sheet_names:
                df = xls.parse(sheet_name)
                text.append(f"\n[FEUILLE '{sheet_name}']\n{df.to_string()}")
        return "\n".join(text), metadata

    @staticmethod
    def _extract_docx(file_path: str, metadata: dict) -> Tuple[str, dict]:
        """Extrait le texte de Word avec python-docx"""
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        metadata["paragraphs"] = len(paragraphs)
        return "\n".join(paragraphs), metadata

    @staticmethod
    def _extract_pptx(file_path: str, metadata: dict) -> Tuple[str, dict]:
        """Extrait le texte de PowerPoint avec python-pptx"""
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        metadata["slides"] = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text.append(f"\n[SLIDE {i+1}]\n" + "\n".join(slide_text))
                metadata[f"slide_{i+1}"] = True  # Marquer la slide dans les métadonnées
        return "\n".join(text), metadata

    @staticmethod
    def _extract_csv(file_path: str, metadata: dict) -> Tuple[str, dict]:
        """Extrait le texte de CSV avec pandas"""
        df = pd.read_csv(file_path)
        metadata["rows"] = len(df)
        return df.to_string(), metadata

    @staticmethod
    def _extract_txt(file_path: str, metadata: dict) -> Tuple[str, dict]:
        """Lit un fichier texte brut"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        metadata["lines"] = len(content.splitlines())
        return content, metadata

    @staticmethod
    def process_directory() -> List[Document]:
        """Traite tous les fichiers du dossier DATA_DIR et retourne les chunks"""
        from config import DATA_DIR
        
        documents = []
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP
        )
        
        # Liste des extensions supportées
        supported_extensions = ['.pdf', '.xlsx', '.xls', '.docx', '.pptx', '.csv', '.txt']
        
        # Parcourir récursivement tous les fichiers
        for root, _, files in os.walk(DATA_DIR):
            for file in files:
                file_path = os.path.join(root, file)
                ext = os.path.splitext(file)[1].lower()
                
                if ext in supported_extensions:
                    logger.info(f"Traitement du fichier: {file_path}")
                    content, metadata = UniversalExtractor.extract_text(file_path)
                    
                    if content:
                        # Ajouter des métadonnées supplémentaires
                        metadata["file_path"] = file_path
                        metadata["file_name"] = os.path.basename(file_path)
                        
                        # Découper le texte
                        chunks = text_splitter.split_text(content)
                        
                        # Créer des documents avec métadonnées
                        for i, chunk in enumerate(chunks):
                            chunk_meta = metadata.copy()
                            chunk_meta["chunk"] = i + 1
                            
                            # Ajouter le numéro de slide pour les PPTX
                            if ext == '.pptx' and 'SLIDE ' in chunk:
                                try:
                                    slide_num = chunk.split('SLIDE ')[1].split(']')[0]
                                    chunk_meta["slide"] = slide_num
                                except:
                                    pass
                            
                            documents.append(Document(page_content=chunk, metadata=chunk_meta))
                    else:
                        logger.warning(f"Aucun contenu extrait pour {file_path}")
                else:
                    logger.warning(f"Format non supporté: {file_path}")
        
        logger.info(f"Nombre total de documents traités: {len(documents)}")
        return documents