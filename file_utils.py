# utils/file_utils.py
import os
from io import BytesIO
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
import pandas as pd # Import pandas
from pptx import Presentation # Import Presentation for PPTX
import logging

logger = logging.getLogger(__name__)

def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Erreur lors de l'extraction du texte du fichier TXT '{file_path}': {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = DocxDocument(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Erreur lors de l'extraction du texte du fichier DOCX '{file_path}': {e}")
        return ""

def extract_text_from_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        full_text = []
        for page in reader.pages:
            full_text.append(page.extract_text() or '') # Handle potentially empty pages
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Erreur lors de l'extraction du texte du fichier PDF '{file_path}': {e}")
        return ""

def extract_text_from_excel(file_path: str) -> str:
    try:
        # Read all sheets into a dictionary of DataFrames
        xls = pd.ExcelFile(file_path)
        all_content = []
        
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            
            # Add sheet name for context
            all_content.append(f"\n--- Contenu de la feuille: {sheet_name} ---\n")
            
            # Iterate through rows and format them explicitly
            # You can customize this formatting based on your typical Excel structure
            for index, row in df.iterrows():
                row_str_parts = []
                for col_name, value in row.items():
                    # Handle NaN values for cleaner output
                    if pd.isna(value):
                        value_str = "N/A"
                    else:
                        value_str = str(value)
                    row_str_parts.append(f"{col_name}: {value_str}")
                all_content.append(f"Ligne {index + 1}: " + ", ".join(row_str_parts))
            all_content.append("\n") # Add a separator between sheets for clarity
            
        return "\n".join(all_content)
    except Exception as e:
        logger.error(f"Erreur lors de l'extraction du texte du fichier Excel '{file_path}': {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Erreur lors de l'extraction du texte du fichier PPTX '{file_path}': {e}")
        return ""